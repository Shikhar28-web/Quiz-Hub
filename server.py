import os
import json
import uuid
import sqlite3
import datetime
from typing import Dict, List, Any, Optional

from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS

# ---------- Configuration ----------
APP_PORT = int(os.environ.get("PORT", 5000))
DATA_DIR = os.path.join(os.getcwd(), "data")
UPLOAD_DIR = os.path.join(DATA_DIR, "uploads")
DB_PATH = os.path.join(DATA_DIR, "quize.db")
STATIC_DIR = os.path.join(os.getcwd(), "quize")

os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(STATIC_DIR, exist_ok=True)

ALLOWED_EXTENSIONS = {"txt", "pdf", "docx", "pptx"}

# ---------- App ----------
app = Flask(__name__, static_folder=STATIC_DIR, static_url_path="/quize")
CORS(app)

# ---------- Database ----------

def get_db_connection() -> sqlite3.Connection:
	conn = sqlite3.connect(DB_PATH)
	conn.row_factory = sqlite3.Row
	return conn


def init_db() -> None:
	conn = get_db_connection()
	cur = conn.cursor()
	cur.execute(
		"""
		CREATE TABLE IF NOT EXISTS tests (
			id TEXT PRIMARY KEY,
			title TEXT,
			created_at TEXT,
			settings_json TEXT,
			questions_json TEXT,
			source_text TEXT
		);
		"""
	)
	cur.execute(
		"""
		CREATE TABLE IF NOT EXISTS submissions (
			id TEXT PRIMARY KEY,
			test_id TEXT,
			submitted_at TEXT,
			student_name TEXT,
			answers_json TEXT,
			score REAL,
			accuracy REAL,
			details_json TEXT
		);
		"""
	)
	cur.execute(
		"""
		CREATE TABLE IF NOT EXISTS ratings (
			id TEXT PRIMARY KEY,
			test_id TEXT,
			question_idx INTEGER,
			rating INTEGER,
			created_at TEXT
		);
		"""
	)
	conn.commit()
	conn.close()


init_db()

# ---------- Utilities ----------

def allowed_file(filename: str) -> bool:
	return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


# ---------- Text Extraction ----------

def extract_text_from_file(file_path: str) -> str:
	ext = ""
	extension = file_path.rsplit(".", 1)[-1].lower()
	if extension == "txt":
		with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
			text = f.read()
	elif extension == "pdf":
		try:
			from PyPDF2 import PdfReader
			reader = PdfReader(file_path)
			pages_text = []
			for page in reader.pages:
				pages_text.append(page.extract_text() or "")
			text = "\n".join(pages_text)
		except Exception:
			text = ""
	elif extension == "docx":
		try:
			import docx  # python-docx
			document = docx.Document(file_path)
			text = "\n".join(p.text for p in document.paragraphs)
		except Exception:
			text = ""
	elif extension == "pptx":
		try:
			from pptx import Presentation
			prs = Presentation(file_path)
			chunks = []
			for slide in prs.slides:
				for shape in slide.shapes:
					if hasattr(shape, "text"):
						chunks.append(shape.text)
			text = "\n".join(chunks)
		except Exception:
			text = ""
	return text


def extract_text_from_url(url: str) -> str:
	try:
		import requests
		from bs4 import BeautifulSoup
		resp = requests.get(url, timeout=20)
		resp.raise_for_status()
		soup = BeautifulSoup(resp.text, "html.parser")
		# Simple visible text extraction
		for script in soup(["script", "style", "noscript"]):
			script.extract()
		text = " ".join(soup.get_text(separator=" ").split())
		return text
	except Exception:
		return ""


# ---------- LLM Integration ----------

class LLMProvider:
	def __init__(self) -> None:
		self.provider = os.environ.get("PREFERRED_LLM", "mock").lower()
		self._gemini = None
		self._hf_key = os.environ.get("HUGGINGFACE_API_KEY")
		self._hf_model = os.environ.get("HF_MODEL", "google/flan-t5-large")
		if self.provider == "gemini":
			try:
				import google.generativeai as genai
				api_key = os.environ.get("GEMINI_API_KEY")
				if not api_key:
					raise RuntimeError("GEMINI_API_KEY not set")
				genai.configure(api_key=api_key)
				self._gemini = genai.GenerativeModel(os.environ.get("GEMINI_MODEL", "gemini-1.5-flash"))
			except Exception:
				self.provider = "mock"

	def _call_gemini(self, prompt: str) -> str:
		assert self._gemini is not None
		try:
			resp = self._gemini.generate_content(prompt)
			return getattr(resp, "text", "") or ""
		except Exception:
			# On any error (e.g., rate limit), return empty to trigger fallback
			return ""

	def _call_huggingface(self, prompt: str) -> str:
		try:
			import requests
			if not self._hf_key:
				return ""
			model = self._hf_model
			headers = {"Authorization": f"Bearer {self._hf_key}"}
			payload = {"inputs": prompt, "parameters": {"max_new_tokens": 800}}
			resp = requests.post(
				f"https://api-inference.huggingface.co/models/{model}",
				headers=headers,
				json=payload,
				timeout=60,
			)
			resp.raise_for_status()
			data = resp.json()
			# HF responses vary; handle both list and dict
			if isinstance(data, list) and data and "generated_text" in data[0]:
				return data[0]["generated_text"]
			if isinstance(data, dict) and "generated_text" in data:
				return data["generated_text"]
			# Some models return {"outputs": [{"text": "..."}]}
			if isinstance(data, dict) and "outputs" in data and data["outputs"]:
				return data["outputs"][0].get("text", "")
			return ""
		except Exception:
			return ""

	def generate_questions(self, source_text: str, settings: Dict[str, Any]) -> List[Dict[str, Any]]:
		"""
		Return list of question dicts with keys:
		- id, type (mcq|multiple_correct|true_false|fill_blank), difficulty (easy|medium|hard),
		- prompt, options (for choices), correct (index or list of indices or string), explanation, topic
		"""
		prompt = self._build_generation_prompt(source_text, settings)
		if self.provider == "gemini":
			raw = self._call_gemini(prompt)
			parsed = self._parse_llm_output(raw)
			if parsed:
				return parsed
		elif self._hf_key:
			raw = self._call_huggingface(prompt)
			parsed = self._parse_llm_output(raw)
			if parsed:
				return parsed
		# Fallback heuristic
		return self._heuristic_generate(source_text, settings)

	def _build_generation_prompt(self, source_text: str, settings: Dict[str, Any]) -> str:
		return (
			"You are an expert educational question setter. Given the source material, generate a JSON array of questions.\n"
			"Each question must be an object with keys: id, type in ['mcq','multiple_correct','true_false','fill_blank'], difficulty in ['easy','medium','hard'], "
			"prompt, options (array for choice types), correct (index for mcq, array of indices for multiple_correct, boolean for true_false, or string for fill_blank), explanation, topic.\n"
			"Adhere to the requested counts per type and difficulty. Keep prompts concise and unambiguous.\n"
			f"Settings JSON: {json.dumps(settings)}\n"
			"Source material begins:\n" + source_text[:20000] + "\nSource material ends."
		)

	def _parse_llm_output(self, raw: str) -> List[Dict[str, Any]]:
		if not raw:
			return []
		# Try to find a JSON array in the output
		start = raw.find("[")
		end = raw.rfind("]")
		if start != -1 and end != -1 and end > start:
			candidate = raw[start : end + 1]
			try:
				data = json.loads(candidate)
				# Basic validation
				valid = []
				for i, q in enumerate(data):
					if isinstance(q, dict) and "prompt" in q and "type" in q:
						q.setdefault("id", str(uuid.uuid4()))
						q.setdefault("difficulty", "medium")
						q.setdefault("topic", "General")
						q.setdefault("explanation", "")
						valid.append(q)
					if len(valid) >= 1:
						continue
					# else skip
				return valid
			except Exception:
				return []
		return []

	def _heuristic_generate(self, source_text: str, settings: Dict[str, Any]) -> List[Dict[str, Any]]:
		import re
		from itertools import islice
		# Very simple generator: split into sentences, create TF and fill blanks, some MCQs from repeated words
		num_questions = int(settings.get("num_questions", 10))
		type_dist = settings.get("type_distribution", {}) or {}
		difficulty_dist = settings.get("difficulty_distribution", {}) or {}
		# Default distributions
		if not type_dist:
			type_dist = {"mcq": max(1, num_questions // 3), "true_false": num_questions // 3, "multiple_correct": num_questions // 6, "fill_blank": num_questions - (num_questions // 3) - (num_questions // 3) - (num_questions // 6)}
		if not difficulty_dist:
			difficulty_dist = {"easy": num_questions // 2, "medium": num_questions // 3, "hard": num_questions - (num_questions // 2) - (num_questions // 3)}

		def cycle_difficulty(idx: int) -> str:
			order = ["easy", "medium", "hard"]
			return order[idx % 3]

		# Prepare sentences
		sentences = re.split(r"(?<=[.!?])\s+", source_text)
		sentences = [s.strip() for s in sentences if len(s.strip()) > 40]
		if not sentences:
			sentences = [source_text[:200]] if source_text else ["Sample context for question generation."]

		questions: List[Dict[str, Any]] = []

		# Fill in the blank
		fill_needed = type_dist.get("fill_blank", 0)
		for idx, s in enumerate(islice(sentences, fill_needed)):
			words = [w for w in re.findall(r"[A-Za-z]{4,}", s)]
			if not words:
				continue
			blank = words[len(words) // 2]
			prompt = s.replace(blank, "____", 1)
			questions.append({
				"id": str(uuid.uuid4()),
				"type": "fill_blank",
				"difficulty": cycle_difficulty(len(questions)),
				"prompt": prompt,
				"correct": blank,
				"explanation": f"The missing word is '{blank}'.",
				"topic": "Key term"
			})

		# True/False
		tf_needed = type_dist.get("true_false", 0)
		for idx, s in enumerate(islice(sentences, tf_needed)):
			truth = (idx % 2 == 0)
			prompt = s if truth else s.replace(" is ", " is not ") if " is " in s else ("Not true: " + s)
			questions.append({
				"id": str(uuid.uuid4()),
				"type": "true_false",
				"difficulty": cycle_difficulty(len(questions)),
				"prompt": prompt,
				"correct": truth,
				"explanation": "Based on the provided source text.",
				"topic": "Comprehension"
			})

		# MCQ and Multiple Correct from keywords
		mcq_needed = type_dist.get("mcq", 0)
		multi_needed = type_dist.get("multiple_correct", 0)
		def top_words(text: str, k: int = 20) -> List[str]:
			freq: Dict[str, int] = {}
			for w in re.findall(r"[A-Za-z]{5,}", text.lower()):
				freq[w] = freq.get(w, 0) + 1
			return [w for w, _ in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:k]]

		keywords = top_words(source_text, 40)
		distractors = [w for w in keywords]

		for idx in range(mcq_needed):
			if not keywords:
				break
			answer = keywords[idx % len(keywords)]
			opts = list({answer} | set(distractors[idx: idx + 5]))[:4]
			if answer not in opts:
				opts[0] = answer
			correct_index = opts.index(answer)
			prompt = f"Which of the following best relates to the topic: '{answer}'?"
			questions.append({
				"id": str(uuid.uuid4()),
				"type": "mcq",
				"difficulty": cycle_difficulty(len(questions)),
				"prompt": prompt,
				"options": opts,
				"correct": correct_index,
				"explanation": f"'{answer}' appears frequently in the source material.",
				"topic": answer.title()
			})

		for idx in range(multi_needed):
			if len(keywords) < 2:
				break
			answers = [keywords[(idx + j) % len(keywords)] for j in range(2)]
			opts = list({*answers} | set(distractors[idx: idx + 6]))[:5]
			correct_indices = [opts.index(a) for a in answers if a in opts]
			prompt = "Select all options that are key topics in the material."
			questions.append({
				"id": str(uuid.uuid4()),
				"type": "multiple_correct",
				"difficulty": cycle_difficulty(len(questions)),
				"prompt": prompt,
				"options": opts,
				"correct": correct_indices,
				"explanation": "Multiple core terms are present.",
				"topic": ", ".join(a.title() for a in answers)
			})

		# Trim to requested count
		return questions[:num_questions]


llm = LLMProvider()

# ---------- Evaluation ----------

def evaluate_answers(questions: List[Dict[str, Any]], answers: Dict[str, Any]) -> Dict[str, Any]:
	total = len(questions)
	correct_count = 0
	details = []
	for idx, q in enumerate(questions):
		qid = q.get("id", str(idx))
		student_answer = answers.get(qid)
		result = {"id": qid, "type": q.get("type"), "prompt": q.get("prompt"), "correct": False}

		q_type = q.get("type")
		if q_type == "mcq":
			correct_index = q.get("correct")
			result["correct_answer"] = correct_index
			result["student_answer"] = student_answer
			if isinstance(student_answer, int):
				if student_answer == correct_index:
					result["correct"] = True

		elif q_type == "multiple_correct":
			correct_indices = set(q.get("correct") or [])
			student_indices = set(student_answer) if isinstance(student_answer, list) else set()
			result["correct_answer"] = list(correct_indices)
			result["student_answer"] = list(student_indices)
			# Only mark correct if the student answered and it exactly matches
			if student_indices and student_indices == correct_indices:
				result["correct"] = True

		elif q_type == "true_false":
			correct_bool = bool(q.get("correct"))
			result["correct_answer"] = correct_bool
			# Only compare if the student explicitly chose True or False
			if isinstance(student_answer, bool):
				result["student_answer"] = student_answer
				if student_answer == correct_bool:
					result["correct"] = True
			else:
				result["student_answer"] = None

		elif q_type == "fill_blank":
			correct_text = str(q.get("correct") or "").strip().lower()
			student_text_raw = student_answer if isinstance(student_answer, str) else ""
			student_text = student_text_raw.strip().lower()
			result["correct_answer"] = q.get("correct")
			result["student_answer"] = student_text_raw
			if student_text and correct_text and student_text == correct_text:
				result["correct"] = True

		else:
			# Unknown type: never mark correct unless explicitly matched
			result["student_answer"] = student_answer
			result["correct_answer"] = q.get("correct")

		if result["correct"]:
			correct_count += 1
		result["explanation"] = q.get("explanation", "")
		details.append(result)
	accuracy = (correct_count / total * 100.0) if total else 0.0
	return {"total": total, "correct": correct_count, "score": correct_count, "accuracy": accuracy, "details": details}


# ---------- Routes ----------

@app.route("/")
def index() -> Any:
	return send_from_directory(STATIC_DIR, "index.html")


@app.route("/test/<test_id>")
def test_page(test_id: str) -> Any:
	return send_from_directory(STATIC_DIR, "index.html")


@app.route("/api/create_test", methods=["POST"])
def create_test() -> Any:
	# Accept multipart (file upload) or JSON (url/text)
	title = request.form.get("title") or (request.json or {}).get("title") or "Untitled Test"
	settings_raw = request.form.get("settings")
	if settings_raw is None and request.is_json:
		settings = (request.json or {}).get("settings", {})
	else:
		settings = json.loads(settings_raw) if settings_raw else {}

	source_text = ""
	# File upload
	if "file" in request.files:
		file = request.files["file"]
		if file and allowed_file(file.filename):
			filename = f"{uuid.uuid4()}_{file.filename}"
			path = os.path.join(UPLOAD_DIR, filename)
			file.save(path)
			source_text = extract_text_from_file(path)
	# URL
	if not source_text:
		url = (request.form.get("url") if not request.is_json else (request.json or {}).get("url"))
		if url:
			source_text = extract_text_from_url(url)
	# Raw text
	if not source_text:
		text = (request.form.get("text") if not request.is_json else (request.json or {}).get("text"))
		if text:
			source_text = text

	if not source_text:
		return jsonify({"error": "No content found from file/url/text."}), 400

	# Generate questions
	questions = llm.generate_questions(source_text, settings)
	if not questions:
		# Ensure we never 500 due to empty LLM result
		questions = llm._heuristic_generate(source_text, settings)
	# Assign missing per-type fields
	for q in questions:
		q.setdefault("id", str(uuid.uuid4()))
		q.setdefault("difficulty", "medium")
		q.setdefault("topic", "General")

	# Persist
	test_id = str(uuid.uuid4())
	now = datetime.datetime.utcnow().isoformat()
	conn = get_db_connection()
	cur = conn.cursor()
	cur.execute(
		"INSERT INTO tests (id, title, created_at, settings_json, questions_json, source_text) VALUES (?, ?, ?, ?, ?, ?)",
		(
			test_id,
			title,
			now,
			json.dumps(settings),
			json.dumps(questions),
			source_text,
		),
	)
	conn.commit()
	conn.close()

	return jsonify({"test_id": test_id, "title": title, "questions": questions})


@app.route("/api/test/<test_id>", methods=["GET"])
def get_test(test_id: str) -> Any:
	conn = get_db_connection()
	cur = conn.cursor()
	cur.execute("SELECT id, title, questions_json FROM tests WHERE id = ?", (test_id,))
	row = cur.fetchone()
	conn.close()
	if not row:
		return jsonify({"error": "Test not found"}), 404
	questions = json.loads(row["questions_json"]) if row["questions_json"] else []
	# Mask correct answers for student delivery; include alongside for evaluation if needed on server
	masked = []
	for q in questions:
		m = {k: v for k, v in q.items() if k != "correct"}
		masked.append(m)
	return jsonify({"test_id": row["id"], "title": row["title"], "questions": masked})


@app.route("/api/submit/<test_id>", methods=["POST"])
def submit_answers(test_id: str) -> Any:
	payload = request.get_json(force=True) or {}
	student_name = payload.get("student_name") or "Anonymous"
	answers = payload.get("answers") or {}
	conn = get_db_connection()
	cur = conn.cursor()
	cur.execute("SELECT questions_json FROM tests WHERE id = ?", (test_id,))
	row = cur.fetchone()
	if not row:
		conn.close()
		return jsonify({"error": "Test not found"}), 404
	questions = json.loads(row["questions_json"]) if row["questions_json"] else []
	result = evaluate_answers(questions, answers)
	submission_id = str(uuid.uuid4())
	now = datetime.datetime.utcnow().isoformat()
	cur.execute(
		"INSERT INTO submissions (id, test_id, submitted_at, student_name, answers_json, score, accuracy, details_json) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
		(
			submission_id,
			test_id,
			now,
			student_name,
			json.dumps(answers),
			float(result.get("score", 0)),
			float(result.get("accuracy", 0.0)),
			json.dumps(result.get("details", [])),
		),
	)
	conn.commit()
	conn.close()
	return jsonify({"submission_id": submission_id, **result})


@app.route("/api/rate/<test_id>", methods=["POST"])
def rate_question(test_id: str) -> Any:
	payload = request.get_json(force=True) or {}
	question_idx = int(payload.get("question_idx", -1))
	rating = int(payload.get("rating", 0))
	if question_idx < 0 or rating < 1 or rating > 5:
		return jsonify({"error": "Invalid rating payload"}), 400
	conn = get_db_connection()
	cur = conn.cursor()
	cur.execute("SELECT 1 FROM tests WHERE id = ?", (test_id,))
	row = cur.fetchone()
	if not row:
		conn.close()
		return jsonify({"error": "Test not found"}), 404
	rating_id = str(uuid.uuid4())
	now = datetime.datetime.utcnow().isoformat()
	cur.execute(
		"INSERT INTO ratings (id, test_id, question_idx, rating, created_at) VALUES (?, ?, ?, ?, ?)",
		(rating_id, test_id, question_idx, rating, now),
	)
	conn.commit()
	conn.close()
	return jsonify({"ok": True})


@app.route("/api/export/<test_id>", methods=["GET"])
def export_test(test_id: str) -> Any:
	# For dev: just return the local link
	base = request.host_url.rstrip("/")
	link = f"{base}/test/{test_id}"
	return jsonify({"link": link})


@app.route("/api/chat", methods=["POST"])
def chat() -> Any:
	payload = request.get_json(force=True) or {}
	message = payload.get("message") or ""
	test_id = payload.get("test_id")
	context = ""
	if test_id:
		conn = get_db_connection()
		cur = conn.cursor()
		cur.execute("SELECT source_text FROM tests WHERE id = ?", (test_id,))
		row = cur.fetchone()
		conn.close()
		if row:
			context = row["source_text"] or ""
	if not message:
		return jsonify({"reply": "Please enter a question."})
	if llm.provider == "mock":
		# Simple echo with hint
		reply = "I'm here to help. " + ("Context: " + context[:200] + " ...\n" if context else "") + f"Your question was: '{message}'."
		return jsonify({"reply": reply})
	# Use LLM
	prompt = (
		"You are a helpful educational assistant. Answer clearly and briefly.\n" +
		("Use this context if relevant:\n" + context[:4000] + "\n" if context else "") +
		"User: " + message + "\nAssistant:"
	)
	if llm.provider == "gemini":
		reply = llm._call_gemini(prompt)
		if not reply:
			reply = "I'm temporarily rate limited. Here is a brief tip based on your context: " + (context[:200] if context else "Try rephrasing your question.")
		return jsonify({"reply": reply.strip()})
	else:
		reply = llm._call_huggingface(prompt)
		if not reply:
			reply = "I'm temporarily unavailable. Please try again shortly."
		return jsonify({"reply": reply.strip()})


# Static files (CSS/JS)
@app.route("/quize/<path:path>")
def send_static(path: str) -> Any:
	return send_from_directory(STATIC_DIR, path)


if __name__ == "__main__":
	app.run(host="0.0.0.0", port=APP_PORT, debug=True) 